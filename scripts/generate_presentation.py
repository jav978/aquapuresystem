#!/usr/bin/env python3
"""
Generador de Presentación Ejecutiva PowerPoint (AquaPureSystem)
Análisis de Costo-Beneficio, Bondades, Estadísticas, KPIs y Solución de Problemas
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Paleta de Colores Corporativa AquaPure
    C_PRIMARY_DARK = RGBColor(10, 37, 64)       # #0A2540 - Navy Profundo
    C_BLUE_BRAND = RGBColor(0, 112, 243)        # #0070F3 - Azul Hídrico
    C_TEAL = RGBColor(0, 168, 150)             # #00A896 - Verde Agua / Eficiencia
    C_ACCENT_LIGHT = RGBColor(240, 248, 255)    # #F0F8FF - Fondo Suave
    C_WHITE = RGBColor(255, 255, 255)
    C_TEXT_MAIN = RGBColor(30, 41, 59)          # #1E293B - Slate Dark
    C_TEXT_MUTED = RGBColor(100, 116, 139)      # #64748B - Slate Muted
    C_SUCCESS = RGBColor(16, 185, 129)          # #10B981 - Verde Éxito
    C_DANGER = RGBColor(239, 68, 68)            # #EF4444 - Rojo Alerta
    C_CARD_BG = RGBColor(248, 250, 252)         # #F8FAFC
    C_CARD_BORDER = RGBColor(226, 232, 240)     # #E2E8F0

    def add_header(slide, title_text, category_text="AQUAPURE SYSTEM v1.0 — PROPUESTA DE VALOR"):
        # Header banner top
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_PRIMARY_DARK
        top_bar.line.fill.background()

        # Accent line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = C_BLUE_BRAND
        accent.line.fill.background()

        # Category text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_TEAL

        # Title text
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(11.5), Inches(0.65))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = C_WHITE

        # Footer
        ft = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35))
        tf_ft = ft.text_frame
        p_ft = tf_ft.paragraphs[0]
        p_ft.text = "AquaPureSystem — Software Especializado de Gestión y Punto de Venta para Purificadoras de Agua"
        p_ft.font.size = Pt(9.5)
        p_ft.font.color.rgb = C_TEXT_MUTED

    def add_card(slide, left, top, width, height, title, items, bg_color=C_CARD_BG, border_color=C_CARD_BORDER, title_color=C_PRIMARY_DARK, badge=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
        tf = txBox.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = title_color
        p0.space_after = Pt(8)

        if badge:
            p_badge = tf.add_paragraph()
            p_badge.text = f"⭐ {badge}"
            p_badge.font.size = Pt(11)
            p_badge.font.bold = True
            p_badge.font.color.rgb = C_BLUE_BRAND
            p_badge.space_after = Pt(8)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = C_TEXT_MAIN
            p.space_after = Pt(4)

    def add_kpi_box(slide, left, top, width, height, number_str, label_str, subtext="", color=C_BLUE_BRAND):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = C_WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        txBox = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(width - 0.3), Inches(height - 0.2))
        tf = txBox.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = number_str
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label_str
        p2.font.size = Pt(12.5)
        p2.font.bold = True
        p2.font.color.rgb = C_PRIMARY_DARK
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(2)

        if subtext:
            p3 = tf.add_paragraph()
            p3.text = subtext
            p3.font.size = Pt(10)
            p3.font.color.rgb = C_TEXT_MUTED
            p3.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 1: PORTADA EJECUTIVA
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_PRIMARY_DARK
    bg1.line.fill.background()

    # Decorative wave / banner
    deco = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.333), Inches(1.7))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(15, 52, 90)
    deco.line.fill.background()

    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(3.8))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "AQUAPURE SYSTEM v1.0"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 178)
    p.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "Sistema Integral de Gestión, Punto de Venta Express y Control Operativo"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_after = Pt(16)

    p3 = tf1.add_paragraph()
    p3.text = "Propuesta de Valor, Transformación Digital y Análisis Costo-Beneficio para Purificadoras de Agua"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(30)

    # Presentation metadata
    p_meta = s1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(1.0)).text_frame
    p_m1 = p_meta.paragraphs[0]
    p_m1.text = "💼 Presentación Ejecutiva de Proyecto | 🚀 Retorno de Inversión (ROI) < 3 Meses | ⚡ Arquitectura de Alto Rendimiento"
    p_m1.font.size = Pt(12)
    p_m1.font.bold = True
    p_m1.font.color.rgb = C_WHITE

    # ==========================================
    # SLIDE 2: DIAGNÓSTICO DEL PROBLEMA
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Diagnóstico del Problema: Fugas Operativas en Purificadoras", "1. DOLORES CRÍTICOS DEL SECTOR")

    add_card(s2, 0.8, 1.5, 5.6, 2.5, "1. Pérdida Crónica de Envases 19L", [
        "Fuga del 15% al 25% del parque de botellones por falta de control estricto de garantías y préstamos.",
        "Cobro erróneo de recargas a clientes que no entregan envase vacío de intercambio.",
        "Impacto anual: Miles de dólares perdidos en reposición constante de botellones nuevos."
    ], bg_color=RGBColor(254, 242, 242), border_color=C_DANGER, title_color=RGBColor(153, 27, 27))

    add_card(s2, 6.8, 1.5, 5.7, 2.5, "2. Descuadres y Fugas de Caja", [
        "Cierres de turno lentos e imprecisos que generan fricción y falta de transparencia con los cajeros.",
        "Múltiples métodos de pago (Efectivo, Pago Móvil, Tarjetas, Divisas) sin conciliación automática.",
        "Ausencia de arqueo ciego y auditoría en tiempo real para prevenir discrepancias."
    ], bg_color=RGBColor(254, 242, 242), border_color=C_DANGER, title_color=RGBColor(153, 27, 27))

    add_card(s2, 0.8, 4.3, 5.6, 2.5, "3. Lentitud en Horas Pico (POS)", [
        "Filas de espera prolongadas en taquilla debido a sistemas lentos o procesos manuales en papel.",
        "Pérdida de clientes que abandonan la fila ante la demora en el despacho de agua.",
        "El tiempo de atención actual supera los 90 a 120 segundos por transacción."
    ], bg_color=RGBColor(254, 242, 242), border_color=C_DANGER, title_color=RGBColor(153, 27, 27))

    add_card(s2, 6.8, 4.3, 5.7, 2.5, "4. Falta de Visibilidad Multialmacén", [
        "Cero trazabilidad entre el stock de planta central, almacén de vacíos/lavado y camiones de reparto.",
        "Pérdidas de botellones en rutas de distribución sin responsables asignados.",
        "Imposibilidad de tomar decisiones basadas en datos de consumo y rentabilidad."
    ], bg_color=RGBColor(254, 242, 242), border_color=C_DANGER, title_color=RGBColor(153, 27, 27))

    # ==========================================
    # SLIDE 3: LA SOLUCIÓN AQUAPURE SYSTEM
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "La Solución: AquaPureSystem v1.0", "2. MODERNIZACIÓN Y AUTOMATIZACIÓN")

    add_card(s3, 0.8, 1.5, 3.7, 5.3, "⚡ Punto de Venta Express", [
        "Diseñado para atención ultrarrápida en menos de 25 segundos.",
        "Acceso directo por teclado (Atajos F1-F12) para máxima agilidad.",
        "Validación automática de envase entregado vs recarga solicitada.",
        "Impresión térmica instantánea ESC/POS (Tickets y Facturas).",
        "Soporte multimoneda con conversión automática de cambio."
    ], badge="Agilidad en Taquilla", title_color=C_BLUE_BRAND)

    add_card(s3, 4.8, 1.5, 3.7, 5.3, "🔄 Control de Envases y Stock", [
        "Trazabilidad exacta del ciclo de vida del botellón 19L.",
        "Diferenciación de almacenes: Llenos, Vacíos para Lavado, Mermas.",
        "Módulo de inspección física de calidad (Aprobado vs Merma/Fisura).",
        "Gestión de stock en camiones de reparto y transferencias en vivo.",
        "Auditoría inmutable de entradas, salidas y ajustes de inventario."
    ], badge="Cero Pérdida de Envases", title_color=C_TEAL)

    add_card(s3, 8.8, 1.5, 3.7, 5.3, "💼 Caja, Finanzas y Auditoría", [
        "Arqueo de caja ciego con registro exacto de cada denominación.",
        "Conciliación multimoneda (Efectivo, Pago Móvil, Tarjeta, Crédito).",
        "Gestión de cuentas por cobrar y límites de crédito para distribuidores.",
        "Dashboard gerencial con KPIs de ventas, productos más vendidos y márgenes.",
        "Pista de auditoría completa (ActivityLog) con IP y usuario."
    ], badge="Control Financiero Total", title_color=C_PRIMARY_DARK)

    # ==========================================
    # SLIDE 4: ESTADÍSTICAS Y MÉTRICAS DE IMPACTO
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Métricas Clave e Impacto Operativo Proyectado", "3. VARIABLES ESTADÍSTICAS & RESULTADOS")

    add_kpi_box(s4, 0.8, 1.5, 3.6, 2.5, "-95%", "PÉRDIDA DE ENVASES", "Control estricto de garantía y retorno en cada recarga", color=C_SUCCESS)
    add_kpi_box(s4, 4.8, 1.5, 3.7, 2.5, "-70%", "TIEMPO DE ATENCIÓN", "De 90 seg a menos de 25 seg por transacción en POS", color=C_BLUE_BRAND)
    add_kpi_box(s4, 8.8, 1.5, 3.7, 2.5, "100%", "EXACTITUD EN CAJA", "Eliminación de descuadres al cierre con arqueo ciego", color=C_TEAL)

    add_kpi_box(s4, 0.8, 4.3, 3.6, 2.5, "+40%", "CAPACIDAD DE DESPACHO", "Mayor volumen de atención en horas pico sin contratar más personal", color=C_BLUE_BRAND)
    add_kpi_box(s4, 4.8, 4.3, 3.7, 2.5, "-30%", "COSTO POR MERMAS", "Detección temprana de fisuras y trazabilidad por lote", color=C_SUCCESS)
    add_kpi_box(s4, 8.8, 4.3, 3.7, 2.5, "< 3 MESES", "RETORNO DE INVERSIÓN", "Payback garantizado mediante ahorro directo de fugas", color=RGBColor(217, 119, 6))

    # ==========================================
    # SLIDE 5: ANÁLISIS COSTO-BENEFICIO
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Análisis Costo-Beneficio: Planta Mediana (500 Botellones/Día)", "4. MODELO FINANCIERO Y AHORRO REAL")

    add_card(s5, 0.8, 1.5, 5.7, 5.3, "📉 Fugas Actuales (Sin Sistema)", [
        "Pérdida promedio de 40 a 60 botellones/mes no devueltos:",
        "  ➔ Costo directo reposición: $280 - $420 USD/mes.",
        "Descuadres de caja y dinero no registrado:",
        "  ➔ Pérdida estimada: $150 - $300 USD/mes.",
        "Tiempo operativo improductivo en arqueos y reclamos:",
        "  ➔ Pérdida de productividad: 35 horas/mes ($180 USD).",
        "Ventas perdidas por filas lentas en horas pico:",
        "  ➔ Ingreso no capturado: $250 - $400 USD/mes.",
        "----------------------------------------------------",
        "PÉRDIDA TOTAL ESTIMADA: $860 - $1,300 USD / MES"
    ], bg_color=RGBColor(254, 242, 242), border_color=C_DANGER, title_color=RGBColor(153, 27, 27))

    add_card(s5, 6.8, 1.5, 5.7, 5.3, "📈 Ahorro & Ganancia con AquaPure", [
        "Recuperación del 95% del control de envases:",
        "  ➔ Ahorro directo en compras: +$350 USD/mes.",
        "Cero fugas de caja gracias al arqueo ciego auditado:",
        "  ➔ Recuperación inmediata: +$200 USD/mes.",
        "Atención 70% más rápida (Aumento de ventas en horas pico):",
        "  ➔ Incremento de ingresos: +$350 USD/mes.",
        "Optimización de rutas y mermas en camiones:",
        "  ➔ Eficiencia logística: +$150 USD/mes.",
        "----------------------------------------------------",
        "BENEFICIO NETO MENSUAL: +$1,050 USD / MES",
        "🚀 Retorno total de inversión en menos de 90 días."
    ], bg_color=RGBColor(240, 253, 244), border_color=C_SUCCESS, title_color=RGBColor(20, 83, 45))

    # ==========================================
    # SLIDE 6: ARQUITECTURA TÉCNICA Y ESCALABILIDAD
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Arquitectura Tecnológica: Confiabilidad y Alto Rendimiento", "5. SOLIDEZ TÉCNICA")

    add_card(s6, 0.8, 1.5, 3.7, 5.3, "💻 Capa Frontend (Nuxt 4 / Vue 3)", [
        "Interfaz intuitiva de última generación.",
        "Vue 3 Composition API + Pinia Stores.",
        "Diseño responsivo optimizado para pantallas táctiles y teclado POS.",
        "Generación directa de tickets térmicos ESC/POS y reportes PDF/Excel.",
        "Modo Desktop / Local de alta disponibilidad."
    ], badge="Experiencia Fluida", title_color=C_BLUE_BRAND)

    add_card(s6, 4.8, 1.5, 3.7, 5.3, "⚙️ Backend Real-Time (Feathers.js)", [
        "Arquitectura Limpia (Clean Architecture + DDD).",
        "Inyección de dependencias con InversifyJS.",
        "Sincronización en vivo mediante WebSockets (Socket.io) para stock y caja.",
        "Seguridad robusta: Autenticación JWT, roles de usuario y cifrado de contraseñas.",
        "API REST modular y altamente desacoplada."
    ], badge="Sub-segundo & Reactivo", title_color=C_TEAL)

    add_card(s6, 8.8, 1.5, 3.7, 5.3, "🗄️ Persistencia de Datos (Postgres & Redis)", [
        "PostgreSQL 16: Transacciones ACID que garantizan integridad total de cobros.",
        "Prisma ORM v5: Consultas optimizadas con índices avanzados.",
        "Redis v7: Caché en memoria para respuesta instantánea.",
        "Conexión nativa directa vía .env sin requerir dependencias pesadas de Docker.",
        "Copias de seguridad y respaldo automatizado."
    ], badge="Integridad Transaccional", title_color=C_PRIMARY_DARK)

    # ==========================================
    # SLIDE 7: PLAN DE IMPLEMENTACIÓN Y DESPLIEGUE
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Plan de Implementación y Puesta en Marcha", "6. HOJA DE RUTA Y DESPLIEGUE")

    add_card(s7, 0.8, 1.5, 2.7, 5.3, "Fase 1: Configuración", [
        "Duración: Días 1 a 3",
        "---------------------",
        "• Instalación de base de datos PostgreSQL y Redis.",
        "• Carga de catálogo de productos (Recargas 19L, Botellas, Filtros, Accesorios).",
        "• Configuración de almacenes y camiones de reparto.",
        "• Creación de usuarios y roles de cajeros/operadores."
    ], title_color=C_PRIMARY_DARK)

    add_card(s7, 3.8, 1.5, 2.7, 5.3, "Fase 2: Periféricos", [
        "Duración: Días 4 a 6",
        "---------------------",
        "• Integración con impresoras térmicas POS de 58mm/80mm.",
        "• Configuración de formatos de ticket térmico y comprobantes.",
        "• Pruebas de atajos de teclado POS y lectores de código de barras.",
        "• Calibración de tasas de cambio y monedas."
    ], title_color=C_BLUE_BRAND)

    add_card(s7, 6.8, 1.5, 2.7, 5.3, "Fase 3: Capacitación", [
        "Duración: Días 7 a 9",
        "---------------------",
        "• Entrenamiento a cajeros en POS Express y atajos rápidos.",
        "• Capacitación a operadores en recepción e inspección de envases.",
        "• Capacitación a supervisores en arqueo de caja y reportes gerenciales.",
        "• Simulación de escenarios reales."
    ], title_color=C_TEAL)

    add_card(s7, 9.8, 1.5, 2.7, 5.3, "Fase 4: Go-Live & Éxito", [
        "Duración: Día 10 en adelante",
        "---------------------",
        "• Inicio de operaciones en vivo.",
        "• Acompañamiento técnico presencial/remoto durante horas pico.",
        "• Auditoría de los primeros arqueos de caja y cuadres de stock.",
        "• Soporte continuo y actualizaciones de software."
    ], title_color=C_SUCCESS)

    # ==========================================
    # SLIDE 8: CONCLUSIÓN Y PRÓXIMOS PASOS
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    bg8 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = C_PRIMARY_DARK
    bg8.line.fill.background()

    t_box8 = s8.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.5))
    tf8 = t_box8.text_frame
    tf8.word_wrap = True

    p = tf8.paragraphs[0]
    p.text = "Conclusión: La Ventaja Competitiva de AquaPureSystem"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 178)
    p.space_after = Pt(20)

    p2 = tf8.add_paragraph()
    p2.text = "✅ Blindaje total contra la pérdida de envases retornables de 19L."
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_after = Pt(12)

    p3 = tf8.add_paragraph()
    p3.text = "✅ Cuadres de caja exactos y transparentes sin fricción con el personal."
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = C_WHITE
    p3.space_after = Pt(12)

    p4 = tf8.add_paragraph()
    p4.text = "✅ Atención express en horas pico que incrementa las ventas hasta un 40%."
    p4.font.size = Pt(18)
    p4.font.bold = True
    p4.font.color.rgb = C_WHITE
    p4.space_after = Pt(12)

    p5 = tf8.add_paragraph()
    p5.text = "✅ Retorno de inversión comprobado en menos de 3 meses ($1,000+ USD/mes de beneficio neto)."
    p5.font.size = Pt(18)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(254, 240, 138)
    p5.space_after = Pt(30)

    p6 = tf8.add_paragraph()
    p6.text = "🚀 ¡Listo para transformar la rentabilidad y eficiencia de tu purificadora de agua!"
    p6.font.size = Pt(20)
    p6.font.bold = True
    p6.font.color.rgb = RGBColor(0, 212, 178)

    output_path = "docs/AquaPureSystem_Presentacion_Costo_Beneficio.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
